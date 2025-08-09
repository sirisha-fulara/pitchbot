import os
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

GENERATED_FOLDER = "generated_decks"
os.makedirs(GENERATED_FOLDER, exist_ok=True)

def create_ppt(title, pitch_text):
    prs= Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Generated Pitch Deck"

    content_slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide
    content_slide.shapes.title.text = "Pitch"
    content_slide.placeholders[1].text = pitch_text

    filename = f"{GENERATED_FOLDER}/{title.replace(' ', '_')}.pptx"
    prs.save(filename)
    return filename

def create_pdf(title, pitch_text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, title, ln=True, align="C")

    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, pitch_text)

    filename = f"{GENERATED_FOLDER}/{title.replace(' ', '_')}.pdf"
    pdf.output(filename)
    return filename